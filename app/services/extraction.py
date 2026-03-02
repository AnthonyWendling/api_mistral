import io
import re
from pathlib import Path

import openpyxl
import pytesseract
from docx import Document as DocxDocument
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation


def _normalize_extension(filename: str | None, content_type: str | None) -> str:
    ext = ""
    if filename:
        ext = Path(filename).suffix.lower()
    if not ext and content_type:
        mime = {
            "application/pdf": ".pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/jpg": ".jpg",
            "image/gif": ".gif",
        }
        ext = mime.get(content_type.split(";")[0].strip(), "")
    return ext or ".bin"


def _sniff_extension(content: bytes) -> str:
    """Détecte l'extension à partir des magic bytes (quand l'URL ne donne pas de nom de fichier, ex. NocoDB)."""
    if not content or len(content) < 8:
        return ".bin"
    if content[:4] == b"%PDF":
        return ".pdf"
    if content[:2] == b"PK":
        # ZIP (DOCX, XLSX, PPTX)
        if b"word/" in content[:2000] or b"[Content_Types].xml" in content[:2000]:
            return ".docx"
        if b"xl/" in content[:2000] or b"xl/workbook" in content[:2000]:
            return ".xlsx"
        if b"ppt/" in content[:2000] or b"ppt/slides" in content[:2000]:
            return ".pptx"
        return ".docx"
    if content[:8] == b"\x89PNG\r\n\x1a\n":
        return ".png"
    if content[:2] == b"\xff\xd8":
        return ".jpg"
    if content[:6] in (b"GIF87a", b"GIF89a"):
        return ".gif"
    return ".bin"


def extract_text_from_pdf(content: bytes, filename: str = "") -> str:
    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n\n".join(parts) if parts else ""


def extract_text_from_docx(content: bytes, filename: str = "") -> str:
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_xlsx(content: bytes, filename: str = "") -> str:
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = " ".join(str(c) if c is not None else "" for c in row).strip()
            if line:
                parts.append(line)
    wb.close()
    return "\n".join(parts) if parts else ""


def extract_text_from_pptx(content: bytes, filename: str = "") -> str:
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n\n".join(parts) if parts else ""


def extract_text_from_image(content: bytes, filename: str = "") -> str:
    img = Image.open(io.BytesIO(content))
    if img.mode not in ("L", "RGB", "RGBA"):
        img = img.convert("RGB")
    return pytesseract.image_to_string(img) or ""


EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".xlsx": extract_text_from_xlsx,
    ".pptx": extract_text_from_pptx,
    ".png": extract_text_from_image,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
    ".gif": extract_text_from_image,
}


def extract_text(content: bytes, filename: str | None = None, content_type: str | None = None) -> str:
    ext = _normalize_extension(filename or "", content_type)
    if ext == ".bin":
        ext = _sniff_extension(content)
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Type de document non supporté: {ext or 'inconnu'}. Supportés: PDF, DOCX, XLSX, PPTX, PNG, JPG, GIF.")
    return extractor(content, filename or "").strip() or ""
